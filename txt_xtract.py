from io import BytesIO, StringIO, TextIOWrapper
from itertools import chain
from PIL import Image
from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.metrics.pairwise import cosine_similarity
from gensim import corpora, models
from gensim.models import KeyedVectors, Word2Vec, ldamodel
from markdown import markdown
from xlrd.sheet import ctype_text

from nltk.tokenize import RegexpTokenizer
from stop_words import get_stop_words
from nltk.stem.snowball import SnowballStemmer
from nltk import ne_chunk, pos_tag, word_tokenize
from nltk.tree import Tree
from nltk.metrics.distance import edit_distance

import PyPDF2
from pptx import Presentation

from xml.dom.minidom import parseString
try:
    from xml.etree.cElementTree import XML
except ImportError:
    from xml.etree.ElementTree import XML

def process_text(file_bytes, ext):

    image = set(['jpg', 'jpeg', 'png', 'tiff', 'tif', 'gif', 'bmp'])
    video = set(['3gp', '3g2', 'avi', 'f4v', 'flv', 'm4v', 'asf', 'wmv', 'mpeg', 'mp4', 'qt'])
    document = set(['txt', 'rtf', 'dotx', 'dot', 'odt', 'pages', 'tex',
                    'pdf', 'ps', 'eps', 'prn', 'md', 'py', 'java', 'scala'])
    open_office = set(['odt', 'ott', 'odm', 'oth', 'ods', 'ots', 'odg',
                       'otg', 'odp', 'otp', 'odf', 'odb', 'odp'])
    doc_x = set(['docx', 'doc'])
    web = set(['html', 'xhtml', 'php', 'js', 'xml', 'war', 'ear' 'dhtml', 'mhtml'])
    spreathseet =  set(['xls', 'xlsx', 'xltx', 'xlt', 'ods', 'xlsb', 'xlsm', 'xltm'])
    presentation = set(['ppt', 'pptx', 'pot', 'potx', 'ppsx',
                        'pps', 'pptm', 'potm', 'ppsm', 'key'])

    text = None

    if ext is None:
        return ""

    if ext == "pdf":
        logging.disable(logging.CRITICAL)
        text = pdf_text_encode(file_bytes, ENCODING)
        logging.disable(logging.NOTSET)
    elif ext == "csv":
        text = csv_text_encode(file_bytes, ENCODING)
    elif ext == "tsv":
        text = csv_text_encode(file_bytes, ENCODING, "\t")
    elif ext == "doc":
        text = doc_text_encode(file_bytes, ENCODING)
    elif ext == "rtf":
        text = None #general_text_extract(file_bytes, ENCODING, "rtf")
    elif ext == "md":
        text = md_text_encode(file_bytes, ENCODING)
    elif ext == "html":
        text = general_text_extract(file_bytes, ENCODING, "html")
    elif ext == "json":
        text = None #general_text_extract(file_bytes, ENCODING, "json")
    elif ext in doc_x:
        text = docx_text_encode(file_bytes, ENCODING)
    elif ext in open_office:
        text = open_office_text_encode(file_bytes, ENCODING)
    elif ext in document:
        text = pure_text_encode(file_bytes, ENCODING)
    elif ext in image:
        text = img_text_encode(file_bytes, ENCODING)
    elif ext in spreathseet:
        text = spreadsheet_text_encode(file_bytes, ENCODING)
    elif ext in presentation:
        text = pptx_text_encode(file_bytes, ENCODING)

    if text is None:
        print("Unprocessing extension: {0}".format(ext))

    if text is not None:
        #print("{0}: {1}".format(ext, text[:100]))
        return text
    else:
        return ""


def pure_text_encode(f_bytes, encoding):
    return f_bytes.decode(encoding)

def csv_text_encode(f_bytes, encoding, delim = ","):
    return f_bytes.decode(encoding).replace(delim, " ")

def pdf_to_text(f_bytes, encoding):

    pdf_reader = PyPDF2.PdfFileReader(BytesIO(pdfFileObj))
    num_pages = pdf_reader.numPages
    while count < num_pages:
        pageObj = pdf_reader.getPage(count)
        count += 1
        text += pageObj.extractText()

    return text

def img_text_encode(f_bytes, encoding):
    im_to_txt = Image.open(BytesIO(f_bytes))
    return pytesseract.image_to_string(img)

def xml_text_encode(f_bytes, encoding):

    doc = parseString(f_bytes)
    paragraphs = doc.getElementsByTagName('text:p')

    text = [str(ch.data) for ch in filter(\
                                     lambda x: x.nodeType == x.TEXT_NODE, \
                                     chain(*[p.childNodes for p in paragraphs]))]
    return " ".join(text)

def open_office_text_encode(f_bytes, encoding):

    """
    Open office files are essentially zipped archives. The key file
    is the content.xml file within the archive, which can then
    be parsed to extract the text.
    """

    open_office_file = zipfile.ZipFile(BytesIO(f_bytes))
    return xml_text_encode(open_office_file.read('content.xml'), encoding)


#This could be improved...
def spreadsheet_text_encode(f_bytes, encoding):

    #UTF-8 is assumed for encoding, which isn't great. May want to modify later.

    wb = xlrd.open_workbook(file_contents = f_bytes)
    text = []
    for sheet in wb.sheets():
        for row in sheet.get_rows():
            filtered_row = filter(lambda x: ctype_text.get(x.ctype, 'not_text') == 'text', row)
            filtered_row = [s.value for s in filtered_row]
            text += [" ".join(filtered_row)]
    return " ".join(text)

#Reference: https://etienned.github.io/posts/extract-text-from-word-docx-simply/
def docx_text_encode(f_bytes, encoding):
    """
    Take the path of a docx file as argument, return the text in unicode.
    """

    WORD_NAMESPACE = '{http://schemas.openxmlformats.org/wordprocessingml/2006/main}'
    PARA = WORD_NAMESPACE + 'p'
    TEXT = WORD_NAMESPACE + 't'

    document = zipfile.ZipFile(BytesIO(f_bytes))
    xml_content = document.read('word/document.xml')
    document.close()
    tree = XML(xml_content)

    paragraphs = []
    for paragraph in tree.getiterator(PARA):
        texts = [node.text
                 for node in paragraph.getiterator(TEXT)
                 if node.text]
        if texts:
            paragraphs.append(''.join(texts))

    return '\n\n'.join(paragraphs)

def general_text_extract(f_bytes, encoding, ext):

    f_type = "_." + ext
    return fulltext.get(BytesIO(f_bytes), name = f_type)

def md_text_encode(f_bytes, encoding):
    html = markdown(f_bytes.decode('utf-8'))
    return fulltext.get(StringIO(html), name = "_.html")

def doc_text_encode(f_bytes, encoding):

    fake_fname = "_.doc"
    with open(fake_fname, "wb") as of:
        of.write(f_bytes)

    text = os.popen("antiword " + fake_fname).read()
    os.system("shred -u " + fake_fname)

    return text

def pptx_text_encode(f_bytes, encoding):

    prs = Presentation(BytesIO(f_bytes))

    text = ""

    for sld in prs.slides:
        for shape in sld.shapes:
            if not shape.has_text_frame:
                continue
            for p in shape.text_frame.paragraphs:
                text += str(p.text) + "\n"

    return text
